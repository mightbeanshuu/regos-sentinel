#!/usr/bin/env python3
"""Build the RegOS Sentinel pitch deck.

Run with /usr/bin/python3 (it is the interpreter that has python-pptx; the
homebrew 3.14 pip is broken on this machine).

    /usr/bin/python3 deck-assets/build_deck.py

Design decisions worth keeping:

* **The palette is the product's palette.** Romer dark — near-black canvas, and
  the four semantic colours the app itself uses: aqua = verified, peach = a
  person is required, periwinkle = computed, coral = a check failed. Lime is
  brand and action ONLY. A judge who watches the film, opens the site and then
  reads the deck should feel one product, not three contractors.

* **Graphics are native PowerPoint shapes, not rasterised SVG.** cairosvg is not
  installed here, but the better reason is that native shapes stay vector, stay
  crisp on a projector, and stay editable by a human before the pitch.

* **The proof slides carry animated GIFs of the real product**, cut from the
  screen recordings at the exact moment the film uses. PowerPoint animates GIFs
  in slideshow mode. Nothing on those slides is a mockup.

* **Every number is real output**, and the ones that are estimates say so. The
  model figures come from app/model/weights.json, regenerated 2026-08-11.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
GIF = ROOT / "gif"
OUT = ROOT.parent / "RegOS_Sentinel_Pitch_2026.pptx"

# ---- Romer dark palette -------------------------------------------------- #
BG          = RGBColor(0x07, 0x07, 0x08)
PANEL       = RGBColor(0x0D, 0x0E, 0x0F)
PANEL_2     = RGBColor(0x14, 0x15, 0x16)
HAIRLINE    = RGBColor(0x24, 0x25, 0x28)
INK         = RGBColor(0xF4, 0xF5, 0xF6)
INK_2       = RGBColor(0xA8, 0xAB, 0xB0)
INK_3       = RGBColor(0x7C, 0x7F, 0x83)
LIME        = RGBColor(0xD0, 0xFE, 0x67)   # brand / action ONLY
AQUA        = RGBColor(0x76, 0xD2, 0xE3)   # verified
PEACH       = RGBColor(0xFF, 0xC2, 0x97)   # a person is required
PERI        = RGBColor(0xC8, 0xCB, 0xFF)   # computed
CORAL       = RGBColor(0xFF, 0x9D, 0x9D)   # a check failed

SANS = "Helvetica Neue"
MONO = "Menlo"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.82)                    # outer margin
CONTENT_W = W - 2 * M

# Product-screenshot slides share one geometry. 7.1in of 16:9 is 3.99in tall,
# so at y=2.62 the image rests 0.17in clear of the footer rule at 6.78.
GIF_Y = Inches(2.62)
GIF_W = Inches(7.1)
GIF_COL = W - Inches(0.82) - (Inches(0.82) + Inches(7.1) + Inches(0.45))


# ---- primitives ----------------------------------------------------------- #
def deck() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def text(s, x, y, w, h, runs, size=18, colour=INK, bold=False, align=PP_ALIGN.LEFT,
         font=SANS, spacing=1.18, anchor=MSO_ANCHOR.TOP, caps=False):
    """`runs` is a string, or a list of (text, {overrides}) for inline colour."""
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    if isinstance(runs, str):
        runs = [(runs, {})]
    for content, over in runs:
        r = p.add_run()
        r.text = content.upper() if caps else content
        f = r.font
        f.name = over.get("font", font)
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.color.rgb = over.get("colour", colour)
        if caps:
            f.size = Pt(over.get("size", size))
    return box


def rule(s, x, y, w, colour=HAIRLINE, thick=Pt(0.75)):
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Emu(9525))
    ln.fill.solid()
    ln.fill.fore_color.rgb = colour
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def panel(s, x, y, w, h, fill=PANEL, border=HAIRLINE, radius=True):
    shape = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    if radius:
        try:
            shape.adjustments[0] = 0.035
        except Exception:
            pass
    return shape


def kicker(s, y, number, label):
    """The editorial section marker: 01 / THE DEFECT."""
    text(s, M, y, Inches(4), Inches(0.3),
         [(f"{number}", {"colour": LIME, "font": MONO, "size": 11, "bold": True}),
          ("   ", {}),
          (label.upper(), {"colour": INK_3, "size": 11, "bold": True})], size=11)


def shield(s, x, y, size=Inches(0.3), colour=LIME):
    """The brand mark, same geometry as the film's and the site's."""
    m = s.shapes.add_shape(MSO_SHAPE.PENTAGON, x, y, size, size)
    m.rotation = 90
    m.fill.background()
    m.line.color.rgb = colour
    m.line.width = Pt(1.5)
    m.shadow.inherit = False
    return m


def footer(s, page):
    rule(s, M, H - Inches(0.72), CONTENT_W)
    text(s, M, H - Inches(0.6), Inches(6), Inches(0.24),
         "RegOS Sentinel · SEBI Securities Market TechSprint @ GFF 2026 · PS2",
         size=9, colour=INK_3)
    text(s, W - M - Inches(1.2), H - Inches(0.6), Inches(1.2), Inches(0.24),
         f"{page:02d}", size=9, colour=INK_3, align=PP_ALIGN.RIGHT, font=MONO)


def headline(s, y, runs, size=44, w=None, lines=1):
    """Sized from the type, not guessed.

    A fixed 1.6in box silently overflowed every headline that wrapped, and the
    overflow landed on whatever was placed beneath it — on the cover, that was
    the sub-paragraph; on the proof slide, the product screenshot.
    """
    text(s, M, y, w or CONTENT_W, Inches(size / 58 * lines + 0.2), runs,
         size=size, bold=True, spacing=1.06)


def gif(s, name, x, y, w):
    """A frame of the real product, animated in slideshow."""
    path = GIF / f"{name}.gif"
    if not path.exists():
        raise SystemExit(f"missing GIF: {path}  (run the ffmpeg step first)")
    frame = panel(s, x - Inches(0.06), y - Inches(0.06),
                  w + Inches(0.12), Inches(0.12) + w * 1080 / 1920,
                  fill=PANEL_2, border=HAIRLINE)
    pic = s.shapes.add_picture(str(path), x, y, width=w)
    return frame, pic


def legend(s, x, y, items):
    """The colour vocabulary, stated once so the deck can rely on it."""
    cx = x
    for label, colour in items:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, y + Inches(0.045),
                                 Inches(0.09), Inches(0.09))
        dot.fill.solid()
        dot.fill.fore_color.rgb = colour
        dot.line.fill.background()
        dot.shadow.inherit = False
        text(s, cx + Inches(0.16), y, Inches(1.9), Inches(0.22),
             label.upper(), size=9, colour=INK_3)
        cx += Inches(0.16) + Inches(1.28)


def stat(s, x, y, w, figure, label, colour=INK, size=38):
    """A figure over its caption.

    The gap is derived from the type size rather than fixed: a 38pt figure needs
    more room beneath it than a 26pt one, and the first cut of this deck set the
    caption 0.6in down for every size, so every large figure sat on top of its
    own label.
    """
    text(s, x, y, w, Inches(size / 52), figure, size=size, bold=True, colour=colour)
    text(s, x, y + Inches(size / 52 + 0.1), w, Inches(0.62), label, size=10.5,
         colour=INK_3, spacing=1.3)


def bar(s, x, y, w, h, frac, colour):
    track = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    track.fill.solid()
    track.fill.fore_color.rgb = PANEL_2
    track.line.fill.background()
    track.shadow.inherit = False
    if frac > 0:
        fill = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Emu(int(w * frac)), h)
        fill.fill.solid()
        fill.fill.fore_color.rgb = colour
        fill.line.fill.background()
        fill.shadow.inherit = False


# ---- slides ---------------------------------------------------------------- #
def s01_cover(prs):
    s = slide(prs)
    shield(s, M, Inches(0.78), Inches(0.34))
    text(s, M + Inches(0.5), Inches(0.8), Inches(6), Inches(0.3),
         "REGOS SENTINEL", size=12, bold=True, colour=INK)
    text(s, M + Inches(2.5), Inches(0.8), Inches(6), Inches(0.3),
         "SEBI TECHSPRINT · PS2 · AGENTIC COMPLIANCE", size=10, colour=INK_3)

    # The headline is bounded well clear of the stat column on the right.
    headline(s, Inches(1.95),
             [("One week.\n", {}), ("From when?", {"colour": LIME})],
             size=72, w=Inches(7.0), lines=2)
    text(s, M, Inches(4.72), Inches(6.7), Inches(1.5),
         "SEBI's cyber framework tells a broker to close a high-severity finding "
         "within one week. It never says one week from when. RegOS Sentinel reads "
         "the circular, finds the gap, and refuses to invent the date.",
         size=14.5, colour=INK_2, spacing=1.45)

    legend(s, M, Inches(6.15), [("computed", PERI), ("verified", AQUA),
                                ("a person decides", PEACH)])

    for i, (fig, lab, col) in enumerate([
        ("205", "pages read,\nof 205", INK),
        ("2,258", "passages\nclassified", INK),
        ("0", "dates\ninvented", LIME),
    ]):
        stat(s, Inches(8.35) + Inches(1.42) * i, Inches(2.3), Inches(1.3),
             fig, lab, col, size=32)
    footer(s, 1)


def s02_defect(prs):
    s = slide(prs)
    kicker(s, Inches(0.8), "01", "The defect")
    headline(s, Inches(1.3), "A period is not a deadline")
    text(s, M, Inches(2.25), Inches(5.6), Inches(1.5),
         "A rule can tell you how long you have and never say when the clock "
         "starts. You cannot compute a date from that, and confidence does not "
         "help. Most tools fill the gap with an assumption.",
         size=14, colour=INK_2, spacing=1.45)

    q = panel(s, M, Inches(3.9), Inches(5.6), Inches(2.15))
    text(s, M + Inches(0.32), Inches(4.15), Inches(4.95), Inches(1.1),
         "“Compensatory controls like virtual patching shall be implemented "
         "for legacy systems for a maximum period of 6 months.”",
         size=13.5, colour=INK, spacing=1.4)
    text(s, M + Inches(0.32), Inches(5.45), Inches(4.95), Inches(0.4),
         "SEBI CSCRF · PR.MA.S3(6) · page 116 — found unprompted, in a document "
         "the model had never seen",
         size=9.5, colour=INK_3, spacing=1.3)

    x2 = M + Inches(6.2)
    for i, (term, detail, col) in enumerate([
        ("Reads the period", "“a maximum period of 6 months” — found.", PERI),
        ("Looks for the trigger", "Nothing in the passage says what starts it.", PERI),
        ("Stops", "No due date is computed. None is guessed.", PEACH),
        ("Hands it over", "The gap goes to a named person, passage attached.", PEACH),
    ]):
        y = Inches(2.3) + Inches(0.92) * i
        chip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x2, y + Inches(0.06),
                                  Inches(0.035), Inches(0.5))
        chip.fill.solid(); chip.fill.fore_color.rgb = col
        chip.line.fill.background(); chip.shadow.inherit = False
        text(s, x2 + Inches(0.24), y, Inches(4.4), Inches(0.3), term, size=15, bold=True)
        text(s, x2 + Inches(0.24), y + Inches(0.32), Inches(4.4), Inches(0.5),
             detail, size=12, colour=INK_2, spacing=1.3)
    footer(s, 2)


def s03_how(prs):
    s = slide(prs)
    kicker(s, Inches(0.8), "02", "How it works")
    headline(s, Inches(1.3), "Agents read. Fixed rules decide. A person judges.")
    text(s, M, Inches(2.3), Inches(9.2), Inches(0.5),
         "Three steps, and one place it deliberately stops.",
         size=14, colour=INK_2)

    steps = [
        ("READ", "The model", PERI,
         "Our own classifier reads every sentence and asks one question: can a "
         "date be computed from this?"),
        ("CHECK", "The fixed rule", AQUA,
         "A deterministic rule reads the same sentence independently. Where they "
         "disagree, the disagreement is the output."),
        ("DECIDE", "The person", PEACH,
         "No date is invented. The gap goes to a named officer, who records a "
         "reading before the tool shows its own."),
    ]
    w = GIF_COL
    for i, (tag, title, col, body) in enumerate(steps):
        x = M + (w + Inches(0.42)) * i
        p = panel(s, x, Inches(3.0), w, Inches(2.75))
        text(s, x + Inches(0.3), Inches(3.28), Inches(2.0), Inches(0.24),
             tag, size=10, bold=True, colour=col)
        text(s, x + Inches(0.3), Inches(3.62), w - Inches(0.6), Inches(0.4),
             title, size=21, bold=True)
        text(s, x + Inches(0.3), Inches(4.18), w - Inches(0.6), Inches(1.4),
             body, size=12, colour=INK_2, spacing=1.45)
        if i < 2:
            text(s, x + w + Inches(0.1), Inches(4.2), Inches(0.3), Inches(0.3),
                 "→", size=15, colour=INK_3)
    footer(s, 3)


def s04_proof(prs):
    s = slide(prs)
    kicker(s, Inches(0.8), "03", "The proof")
    # Kept to one line on purpose: at 32pt this wrapped and left "ago" orphaned
    # on a line of its own, which is the cheapest-looking thing a headline can do.
    headline(s, Inches(1.26),
             [("The same defect, ", {}),
              ("three weeks ago", {"colour": LIME})], size=40)
    gif(s, "upload_case", M, GIF_Y, GIF_W)
    x = M + GIF_W + Inches(0.45)
    text(s, x, GIF_Y, GIF_COL, Inches(2.6),
         "SEBI circular of 23 July 2026, uploaded live. Page 3, paragraph 6:\n\n"
         "“Processing entities shall provide monthly reports … for a "
         "period of 6 months.”\n\n"
         "Six months from what? The circular never says.",
         size=12.5, colour=INK_2, spacing=1.45)
    text(s, x, Inches(5.3), GIF_COL, Inches(0.3),
         "MODEL  ·  PERIOD_ONLY  0.9799", size=11, bold=True, colour=PERI, font=MONO)
    text(s, x, Inches(5.65), GIF_COL, Inches(0.3),
         "FIXED RULE  ·  AGREES", size=11, bold=True, colour=AQUA, font=MONO)
    text(s, x, Inches(6.05), GIF_COL, Inches(0.5),
         "Two paragraphs earlier, the same drafter DID state a trigger — and the "
         "model scores that 0.9943. Same page. Read both and check its work.",
         size=10, colour=INK_3, spacing=1.35)
    footer(s, 4)


def s05_refuses(prs):
    s = slide(prs)
    kicker(s, Inches(0.8), "04", "The part most tools hide")
    headline(s, Inches(1.28), "It would rather say nothing than guess", size=36)
    gif(s, "blocked", M, GIF_Y, GIF_W)
    x = M + GIF_W + Inches(0.45)
    text(s, x, GIF_Y, GIF_COL, Inches(2.2),
         "The period is there. The starting point is not.\n\n"
         "So the due date stays empty, the clock stays blocked, and the gap is "
         "put in front of a person — with the passage attached.",
         size=12.5, colour=INK_2, spacing=1.45)
    p = panel(s, x, Inches(4.85), GIF_COL, Inches(1.5), fill=PANEL_2)
    text(s, x + Inches(0.26), Inches(5.08), Inches(3.1), Inches(1.1),
         "An invented deadline is worse than an admitted gap. One is a wrong "
         "answer a regulator can act on; the other is a question a compliance "
         "officer can answer.",
         size=11.5, colour=PEACH, spacing=1.4)
    footer(s, 5)


def s06_decision(prs):
    s = slide(prs)
    kicker(s, Inches(0.8), "05", "Human in the loop")
    headline(s, Inches(1.28), "A judgement, not a rubber stamp", size=38)
    gif(s, "decide", M, GIF_Y, GIF_W)
    x = M + GIF_W + Inches(0.45)
    for i, (t, d) in enumerate([
        ("The order is enforced",
         "You cannot record a reading until every cited section has been opened "
         "and fingerprinted."),
        ("Reading before reveal",
         "Your own reading is timestamped BEFORE the tool shows you what it thinks."),
        ("A named officer",
         "The decision carries a person, a role, and a written reason."),
    ]):
        y = GIF_Y + Inches(1.15) * i
        text(s, x, y, GIF_COL, Inches(0.3), t, size=14, bold=True, colour=PEACH)
        text(s, x, y + Inches(0.34), GIF_COL, Inches(0.8), d,
             size=11.5, colour=INK_2, spacing=1.4)
    footer(s, 6)


def s07_model(prs):
    """The slide the deck was missing: RegOS trained its own classifier."""
    s = slide(prs)
    kicker(s, Inches(0.8), "06", "Our own model")
    headline(s, Inches(1.28),
             [("We did not call an API. ", {}), ("We trained it.", {"colour": LIME})],
             size=38)
    text(s, M, Inches(2.25), Inches(5.5), Inches(1.25),
         "A multinomial logistic-regression timing classifier, 23 hand-designed "
         "linguistic features, pure Python standard library — no numpy, no "
         "sklearn, no network call. It ships inside the API and runs offline.",
         size=13, colour=INK_2, spacing=1.45)

    for i, (fig, lab, col) in enumerate([
        ("0.842", "accuracy on documents\nit has never seen", INK),
        ("0.946", "recall on PERIOD_ONLY,\nthe defect that matters", LIME),
        ("388", "sentences hand-labelled\nfrom 36 SEBI sources", INK),
    ]):
        stat(s, M + Inches(1.85) * i, Inches(3.7), Inches(1.7), fig, lab, col, size=32)

    text(s, M, Inches(5.35), Inches(5.5), Inches(0.9),
         "Held out a whole document at a time, so every score comes from a model "
         "that had never read a single sentence of that circular. Sentence-level "
         "cross-validation flatters it to 0.892; we quote the harder number.",
         size=10, colour=INK_3, spacing=1.4)

    # The learned weights ARE the argument — show them.
    x = M + Inches(6.35)
    p = panel(s, x, Inches(2.15), Inches(4.85), Inches(4.25))
    text(s, x + Inches(0.32), Inches(2.42), Inches(4.2), Inches(0.3),
         "WHAT IT LEARNED, UNPROMPTED", size=10, bold=True, colour=INK_3)
    weights = [
        ("has_absolute_date", 3.82, "a date is stated", PERI),
        ("periodicity_word", 2.45, "“quarterly”, “annual”", PERI),
        ("duration_without_clock_start", 2.07, "the defect", LIME),
        ("urgency_strong", 2.05, "“immediately”", PEACH),
        ("reported_past", 1.49, "talks about a duty", AQUA),
        ("interrogative", -1.14, "a question, not a duty", AQUA),
    ]
    top = Inches(2.88)
    for i, (name, val, gloss, col) in enumerate(weights):
        y = top + Inches(0.5) * i
        text(s, x + Inches(0.32), y, Inches(2.9), Inches(0.22),
             name, size=9.5, colour=INK, font=MONO)
        text(s, x + Inches(0.32), y + Inches(0.2), Inches(2.9), Inches(0.22),
             gloss, size=8.5, colour=INK_3)
        bar(s, x + Inches(3.35), y + Inches(0.05), Inches(1.0), Inches(0.11),
            min(abs(val) / 3.9, 1.0), col)
        text(s, x + Inches(3.35), y + Inches(0.2), Inches(1.0), Inches(0.22),
             f"{val:+.2f}", size=9, colour=col, font=MONO)
    rule(s, x + Inches(0.32), Inches(5.82), Inches(4.2))
    text(s, x + Inches(0.32), Inches(5.98), Inches(4.2), Inches(0.3),
         "Inspectable, versioned with its training data, and free to run.",
         size=9.5, colour=INK_3)
    footer(s, 7)


def s08_model_vs_rule(prs):
    s = slide(prs)
    kicker(s, Inches(0.8), "07", "Why two readers")
    headline(s, Inches(1.28), "The model never gets the last word", size=38)
    text(s, M, Inches(2.3), Inches(9.4), Inches(0.8),
         "The classifier and a deterministic rule read every passage "
         "independently. Neither overrules the other — where they disagree, the "
         "disagreement IS the output, and it goes to a person.",
         size=14, colour=INK_2, spacing=1.45)

    cards = [
        ("THE MODEL", PERI, "Reads wording it has never seen. Generalises. "
         "Returns a class and a confidence.", "learned from 388 labelled sentences"),
        ("THE FIXED RULE", AQUA, "Pattern-matched, auditable, unchanging. Cannot "
         "be surprised, cannot generalise.", "written down, not trained"),
        ("THE PERSON", PEACH, "Sees both reads and the passage. Records a "
         "decision with a reason.", "the only one who may create a duty"),
    ]
    w = GIF_COL
    for i, (title, col, body, foot) in enumerate(cards):
        x = M + (w + Inches(0.42)) * i
        panel(s, x, Inches(3.55), w, Inches(2.35))
        text(s, x + Inches(0.3), Inches(3.82), Inches(2.6), Inches(0.24),
             title, size=10, bold=True, colour=col)
        text(s, x + Inches(0.3), Inches(4.22), w - Inches(0.6), Inches(1.1),
             body, size=12.5, colour=INK, spacing=1.4)
        text(s, x + Inches(0.3), Inches(5.42), w - Inches(0.6), Inches(0.3),
             foot, size=9.5, colour=INK_3)
    text(s, M, Inches(6.2), Inches(11.5), Inches(0.4),
         "We never train the model on the rule's output. Do that and they can "
         "never disagree — the agreement becomes a tautology and the check "
         "becomes theatre. Every one of the 388 labels was assigned by reading.",
         size=10.5, colour=INK_3, spacing=1.4)
    footer(s, 8)


def s09_assistants(prs):
    s = slide(prs)
    kicker(s, Inches(0.8), "08", "AI assistants")
    headline(s, Inches(1.28), "Four readers, one job each", size=40)
    text(s, M, Inches(2.3), Inches(9.4), Inches(0.4),
         "No assistant holds a tool that writes. The toolbox is the permission model.",
         size=14, colour=INK_2)

    items = [
        ("Reference resolver", "Checks that every “see Table 19” really "
         "points at Table 19.", "5 steps", AQUA),
        ("Source scout", "Spots when SEBI's wording quietly moves.", "11 steps", PERI),
        ("Adversary", "Tries to break our own conclusions before a regulator "
         "can.", "2 steps", PEACH),
        ("Extractor", "Asks of every sentence: can this make a calendar date?",
         "8 steps", PERI),
    ]
    w = Inches(2.68)
    for i, (name, job, steps, col) in enumerate(items):
        x = M + (w + Inches(0.28)) * i
        panel(s, x, Inches(3.1), w, Inches(2.05))
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.26), Inches(3.36),
                                 Inches(0.14), Inches(0.14))
        dot.fill.solid(); dot.fill.fore_color.rgb = col
        dot.line.fill.background(); dot.shadow.inherit = False
        text(s, x + Inches(0.26), Inches(3.62), w - Inches(0.5), Inches(0.3),
             name, size=14.5, bold=True)
        text(s, x + Inches(0.26), Inches(4.0), w - Inches(0.5), Inches(0.85),
             job, size=11, colour=INK_2, spacing=1.35)
        text(s, x + Inches(0.26), Inches(4.82), w - Inches(0.5), Inches(0.24),
             steps, size=9.5, colour=INK_3, font=MONO)

    p = panel(s, M, Inches(5.45), CONTENT_W, Inches(0.95), fill=PANEL_2)
    text(s, M + Inches(0.32), Inches(5.68), Inches(11), Inches(0.55),
         "The adversary found a real defect on its first run: a compiled "
         "obligation cited FAQ Q17(b) for a clock-start that Q17(b) does not "
         "state. Sixty tests had missed it. That regression is now a build gate.",
         size=12, colour=INK, spacing=1.4)
    footer(s, 9)


def s10_record(prs):
    s = slide(prs)
    kicker(s, Inches(0.8), "09", "The record")
    headline(s, Inches(1.28), "A trace you can recompute", size=40)
    gif(s, "record", M, GIF_Y, GIF_W)
    x = M + GIF_W + Inches(0.45)
    text(s, x, GIF_Y, GIF_COL, Inches(2.0),
         "Every step an assistant takes carries the digest of the step before "
         "it. Edit one and verification fails.\n\n"
         "That is what makes the trace evidence rather than a log.",
         size=12.5, colour=INK_2, spacing=1.45)
    for i, (fig, lab) in enumerate([("26", "steps, every one hashed"),
                                    ("153", "tests, all passing"),
                                    ("28", "real SEBI PDFs, 0 problems")]):
        y = Inches(4.65) + Inches(0.66) * i
        text(s, x, y, Inches(0.9), Inches(0.35), fig, size=20, bold=True, colour=LIME,
             font=MONO)
        text(s, x + Inches(0.95), y + Inches(0.08), Inches(2.7), Inches(0.3),
             lab, size=11, colour=INK_2)
    footer(s, 10)


def s11_fit(prs):
    s = slide(prs)
    kicker(s, Inches(0.8), "10", "Why it fits SEBI")
    headline(s, Inches(1.28), "The regulator asked for this, twice", size=38)

    cols = [
        ("CSCRF FAQ · Q30", AQUA,
         "“MIIs and Qualified REs shall build an automated tool and suitable "
         "dashboard … for submitting compliance.”",
         "It is also a scored parameter of the Cyber Capability Index. The "
         "product IS the mandated deliverable — and building it scores a point."),
        ("AI advisory · 5 May 2026", PERI,
         "“All REs need to prepare a long-term plan for usage of AI in "
         "detection and autonomous/agentic mitigation.”",
         "Annexure-A, item 10, under IT Committee guidance. SEBI is asking for "
         "agentic AI with a human gate. That is exactly this architecture."),
    ]
    w = Inches(5.65)
    for i, (src, col, quote, gloss) in enumerate(cols):
        x = M + (w + Inches(0.4)) * i
        panel(s, x, Inches(2.35), w, Inches(3.5))
        text(s, x + Inches(0.34), Inches(2.62), w - Inches(0.7), Inches(0.26),
             src.upper(), size=10, bold=True, colour=col)
        text(s, x + Inches(0.34), Inches(3.05), w - Inches(0.7), Inches(1.3),
             quote, size=15, colour=INK, spacing=1.4)
        rule(s, x + Inches(0.34), Inches(4.5), w - Inches(0.7))
        text(s, x + Inches(0.34), Inches(4.72), w - Inches(0.7), Inches(1.0),
             gloss, size=11.5, colour=INK_2, spacing=1.45)
    text(s, M, Inches(6.2), CONTENT_W, Inches(0.4),
         "Both verified against the published PDFs on sebi.gov.in, not quoted "
         "from memory.", size=10.5, colour=INK_3)
    footer(s, 11)


def s12_honest(prs):
    s = slide(prs)
    kicker(s, Inches(0.8), "11", "What is real, and what is not")
    headline(s, Inches(1.28), "The limitations are on the slide too", size=38)
    text(s, M, Inches(2.28), Inches(9.6), Inches(0.4),
         "A tool that claims certainty it does not have is the thing this product "
         "exists to prevent.", size=13.5, colour=INK_2)

    real = [
        "205-page CSCRF framework: 205/205 pages read, 2,258 passages classified",
        "28 real SEBI PDFs driven end to end — 0 problems",
        "The hero find was unprompted, on a document the model had never seen",
        "153 tests, hash-chain verification, adversary regression gate",
    ]
    limits = [
        "The Cyber Capability Index computes 8 of 23 parameters and abstains on "
        "the rest — it never reports 0 for “unknown”",
        "Weights are ours and labelled as ours; SEBI's 23 parameter names were "
        "not obtainable and are not invented",
        "Uploading the 205-page framework takes ~105s on Render's free tier",
        "Bilingual Hindi passages are marked not-assessed, not silently read",
    ]
    for i, (title, col, rows) in enumerate([("WHAT WORKS", LIME, real),
                                            ("WHAT IT WILL NOT DO", PEACH, limits)]):
        x = M + Inches(6.1) * i
        text(s, x, Inches(3.05), Inches(5.5), Inches(0.26), title, size=10,
             bold=True, colour=col)
        for j, row in enumerate(rows):
            y = Inches(3.45) + Inches(0.72) * j
            d = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.09),
                                   Inches(0.03), Inches(0.03))
            d.fill.solid(); d.fill.fore_color.rgb = col
            d.line.fill.background(); d.shadow.inherit = False
            text(s, x + Inches(0.22), y, Inches(5.25), Inches(0.65), row,
                 size=11.5, colour=INK_2, spacing=1.38)
    footer(s, 12)


def s13_close(prs):
    s = slide(prs)
    shield(s, M, Inches(2.15), Inches(0.5))
    headline(s, Inches(2.95),
             [("Decision support that ", {}), ("shows its work", {"colour": LIME})],
             size=54)
    text(s, M, Inches(4.15), Inches(7.5), Inches(0.5),
         "— and stops where the source stops.", size=20, colour=INK_2)
    rule(s, M, Inches(5.1), Inches(5.2))
    text(s, M, Inches(5.35), Inches(6.5), Inches(0.9),
         "regos-sentinel.vercel.app\ngithub.com/mightbeanshuu/regos-sentinel",
         size=13, colour=INK, font=MONO, spacing=1.5)
    text(s, W - M - Inches(4.6), Inches(5.35), Inches(4.6), Inches(1.0),
         "A prototype built for the SEBI Securities Market TechSprint at GFF "
         "2026. It supports a compliance decision; it is not legal advice, not a "
         "SEBI determination, and nothing is filed automatically.",
         size=9.5, colour=INK_3, align=PP_ALIGN.RIGHT, spacing=1.4)
    footer(s, 13)


def main() -> None:
    prs = deck()
    for build in (s01_cover, s02_defect, s03_how, s04_proof, s05_refuses,
                  s06_decision, s07_model, s08_model_vs_rule, s09_assistants,
                  s10_record, s11_fit, s12_honest, s13_close):
        build(prs)
    prs.save(OUT)
    size = OUT.stat().st_size / 1_000_000
    print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides, {size:.1f} MB)")


if __name__ == "__main__":
    main()
